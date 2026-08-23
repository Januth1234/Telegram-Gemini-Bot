#!/usr/bin/env python3
"""
Orin AI PC Agent v1.0
─────────────────────
Run this script on your Windows/Mac/Linux PC to pair it with orinai.org.
Once paired, tasks dispatched from your phone or any device will execute
here automatically.

Requirements:
  pip install requests pyautogui pillow python-pptx python-docx

Usage:
  python orin-pc-agent.py
  # Enter the 6-char pair code shown on orinai.org → Agent → PC Agent
"""

import hashlib
import hmac
import json
import os
import platform
import shutil
import subprocess
import sys
import tempfile
import time
import traceback
from pathlib import Path
from datetime import datetime

# ── Try to import optional dependencies ──────────────────────────────────────
try:
    import requests
except ImportError:
    print("[Orin] Installing requests...")
    subprocess.run([sys.executable, "-m", "pip", "install", "requests"], check=True)
    import requests

try:
    import pyautogui
    pyautogui.FAILSAFE = True  # move mouse to corner to abort
    HAS_PYAUTOGUI = True
except ImportError:
    HAS_PYAUTOGUI = False
    print("[Orin] Note: pyautogui not installed — type_text will be disabled")
    print("       Install with: pip install pyautogui")

try:
    from PIL import ImageGrab
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("[Orin] Note: python-pptx not installed — create_ppt will be disabled")
    print("       Install with: pip install python-pptx")

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

# ── Configuration ─────────────────────────────────────────────────────────────
BASE_URL   = "https://orinai.org/api/executor"
STATE_FILE = Path.home() / ".orin_agent_state.json"
POLL_INTERVAL_SEC = 1      # brief gap only — the server long-polls (wait=20) until a job is due
LONG_POLL_WAIT = 20       # seconds the server holds /agent/jobs/next open before returning empty
PING_INTERVAL_SEC = 60     # keepalive ping every 60 seconds


# ── HMAC signing ─────────────────────────────────────────────────────────────
def sign_request(secret: str, body: str) -> tuple[str, str]:
    """Returns (timestamp_str, signature_hex)"""
    ts = str(int(time.time()))
    payload = f"{ts}\n{body}"
    sig = hmac.new(secret.encode(), payload.encode(), hashlib.sha256).hexdigest()
    return ts, sig


def agent_post(path: str, pair_id: str, secret: str, data: dict) -> dict:
    body = json.dumps({**data, "pair_id": pair_id}, separators=(",", ":"))
    ts, sig = sign_request(secret, body)
    r = requests.post(
        f"{BASE_URL}{path}",
        data=body,
        headers={
            "Content-Type": "application/json",
            "X-Timestamp": ts,
            "X-Signature": sig,
        },
        timeout=30,
    )
    r.raise_for_status()
    return r.json()


# ── State persistence ─────────────────────────────────────────────────────────
def load_state() -> dict:
    try:
        return json.loads(STATE_FILE.read_text())
    except Exception:
        return {}


def save_state(state: dict):
    STATE_FILE.write_text(json.dumps(state, indent=2))


# ── Pairing ───────────────────────────────────────────────────────────────────
def do_pair(pair_code: str) -> dict:
    """Ask user for pair_id + code, complete handshake, return {pair_id, hmac_secret}"""
    pair_id = input("  Pair ID (copy from orinai.org → PC Agent → pair_id shown below code): ").strip()
    r = requests.post(
        f"{BASE_URL}/pair/agent-handshake",
        json={"pair_id": pair_id, "pair_code": pair_code.upper()},
        timeout=30,
    )
    if r.status_code != 200:
        raise Exception(f"Handshake failed: {r.json().get('error', r.text)}")
    data = r.json()
    return {"pair_id": pair_id, "hmac_secret": data["hmac_secret"]}


# ── Task executors ────────────────────────────────────────────────────────────

def exec_run_command(params: dict) -> dict:
    cmd = params.get("command", "")
    if not cmd:
        return {"error": "No command provided"}
    try:
        result = subprocess.run(
            cmd, shell=True, capture_output=True, text=True, timeout=30
        )
        return {
            "stdout": result.stdout[:2000],
            "stderr": result.stderr[:500],
            "returncode": result.returncode,
            "message": f"Command exited with code {result.returncode}",
        }
    except subprocess.TimeoutExpired:
        return {"error": "Command timed out after 30 seconds"}


def exec_open_app(params: dict) -> dict:
    app = params.get("app", "")
    if not app:
        return {"error": "No app specified"}
    system = platform.system()
    try:
        if system == "Windows":
            os.startfile(app)
        elif system == "Darwin":  # macOS
            subprocess.Popen(["open", "-a", app])
        else:  # Linux
            subprocess.Popen([app])
        return {"message": f"Launched: {app}"}
    except Exception as e:
        # Try as a shell command as fallback
        try:
            subprocess.Popen(app, shell=True)
            return {"message": f"Launched (shell): {app}"}
        except Exception as e2:
            return {"error": str(e2)}


def exec_type_text(params: dict) -> dict:
    if not HAS_PYAUTOGUI:
        return {"error": "pyautogui not installed. Run: pip install pyautogui"}
    text = params.get("text", "")
    if not text:
        return {"error": "No text to type"}
    time.sleep(1)  # give user a moment to focus target window
    pyautogui.typewrite(text, interval=0.03)
    return {"message": f"Typed {len(text)} characters"}


def exec_screenshot(params: dict) -> dict:
    filename = params.get("filename", f"screenshot-{int(time.time())}.png")
    save_path = Path.home() / "Desktop" / filename
    try:
        if HAS_PYAUTOGUI:
            img = pyautogui.screenshot()
            img.save(str(save_path))
        elif HAS_PIL:
            img = ImageGrab.grab()
            img.save(str(save_path))
        else:
            # System fallback
            if platform.system() == "Windows":
                subprocess.run(["snippingtool", "/clip"], timeout=5)
            elif platform.system() == "Darwin":
                subprocess.run(["screencapture", str(save_path)], timeout=5)
            else:
                subprocess.run(["gnome-screenshot", "-f", str(save_path)], timeout=5)
        return {"message": f"Screenshot saved: {save_path}", "path": str(save_path)}
    except Exception as e:
        return {"error": str(e)}


def exec_create_ppt(params: dict) -> dict:
    if not HAS_PPTX:
        return {"error": "python-pptx not installed. Run: pip install python-pptx"}

    title     = params.get("title", "Presentation")
    subtitle  = params.get("subtitle", "Created by Orin AI")
    points    = params.get("points", [])
    direction = params.get("direction", "")

    prs = Presentation()
    prs.slide_width  = 9144000   # 10 inches in EMUs
    prs.slide_height = 5143500   # 5.63 inches

    # ── Slide 1: Title ────────────────────────────────────────────────────
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    if slide.placeholders[1:]:
        slide.placeholders[1].text = subtitle
    # Indigo background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x63, 0x66, 0xF1)

    # ── Slide 2: Key Points ───────────────────────────────────────────────
    if points:
        bullet_layout = prs.slide_layouts[1]
        s2 = prs.slides.add_slide(bullet_layout)
        s2.shapes.title.text = "Key Points"
        tf = s2.placeholders[1].text_frame
        tf.clear()
        for i, pt in enumerate(points):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(pt)
            p.level = 0

    # ── Slide 3: Summary ─────────────────────────────────────────────────
    s3 = prs.slides.add_slide(prs.slide_layouts[1])
    s3.shapes.title.text = "Summary"
    s3.placeholders[1].text = f"Direction: {direction}\n{subtitle}" if direction else subtitle

    # ── Save ──────────────────────────────────────────────────────────────
    safe_name = "".join(c for c in title if c.isalnum() or c in " -_")[:40]
    filename  = f"{safe_name or 'presentation'}.pptx"
    save_path = Path.home() / "Desktop" / filename
    prs.save(str(save_path))

    # Try to open it
    try:
        if platform.system() == "Windows":
            os.startfile(str(save_path))
        elif platform.system() == "Darwin":
            subprocess.Popen(["open", str(save_path)])
        else:
            subprocess.Popen(["xdg-open", str(save_path)])
    except Exception:
        pass

    return {"message": f"Created: {save_path}", "path": str(save_path), "slides": 3}


def exec_create_doc(params: dict) -> dict:
    if not HAS_DOCX:
        return {"error": "python-docx not installed. Run: pip install python-docx"}

    title   = params.get("title", "Document")
    content = params.get("content", "")

    doc = Document()
    doc.add_heading(title, 0)
    doc.add_paragraph(content or "Created by Orin AI")

    safe_name = "".join(c for c in title if c.isalnum() or c in " -_")[:40]
    filename  = f"{safe_name or 'document'}.docx"
    save_path = Path.home() / "Desktop" / filename
    doc.save(str(save_path))

    try:
        if platform.system() == "Windows": os.startfile(str(save_path))
        elif platform.system() == "Darwin": subprocess.Popen(["open", str(save_path)])
        else: subprocess.Popen(["xdg-open", str(save_path)])
    except Exception:
        pass

    return {"message": f"Created: {save_path}", "path": str(save_path)}


def exec_web_search(params: dict) -> dict:
    query = params.get("query", "")
    if not query:
        return {"error": "No query"}
    url = f"https://www.google.com/search?q={requests.utils.quote(query)}"
    try:
        import webbrowser
        webbrowser.open(url)
        return {"message": f"Opened search: {query}", "url": url}
    except Exception as e:
        return {"error": str(e)}


def exec_custom(params: dict) -> dict:
    code = params.get("code", "")
    if not code:
        return {"error": "No code provided"}
    try:
        local_vars: dict = {}
        exec(code, {"__builtins__": __builtins__}, local_vars)
        return {"message": "Code executed", "output": str(local_vars.get("result", ""))}
    except Exception as e:
        return {"error": str(e), "traceback": traceback.format_exc()[:500]}



def exec_spotify(params: dict) -> dict:
    """Control Spotify via spotipy (requires: pip install spotipy)"""
    try:
        import spotipy
        from spotipy.oauth2 import SpotifyOAuth
    except ImportError:
        return {"error": "spotipy not installed. Run: pip install spotipy"}
    
    client_id     = params.get("client_id", os.getenv("SPOTIFY_CLIENT_ID", ""))
    client_secret = params.get("client_secret", os.getenv("SPOTIFY_CLIENT_SECRET", ""))
    redirect_uri  = params.get("redirect_uri", "http://localhost:8888/callback")
    action        = params.get("action", "status")  # play|pause|next|prev|search_play|status|volume
    query         = params.get("query", "")
    volume        = params.get("volume", None)
    
    if not client_id or not client_secret:
        return {"error": "Set SPOTIFY_CLIENT_ID + SPOTIFY_CLIENT_SECRET env vars or pass in params."}
    
    scope = "user-read-playback-state,user-modify-playback-state,user-read-currently-playing,streaming"
    sp = spotipy.Spotify(auth_manager=SpotifyOAuth(
        client_id=client_id, client_secret=client_secret,
        redirect_uri=redirect_uri, scope=scope, cache_path=str(Path.home()/".orin_spotify_cache")
    ))
    
    # Get active device
    devices = sp.devices().get("devices", [])
    active  = next((d for d in devices if d["is_active"]), devices[0] if devices else None)
    dev_id  = active["id"] if active else None
    
    if action == "search_play" and query:
        results = sp.search(q=query, type="track", limit=1)
        items = results["tracks"]["items"]
        if not items: return {"error": f"No track found: {query}"}
        track = items[0]
        sp.start_playback(device_id=dev_id, uris=[track["uri"]])
        return {"message": f"▶ {track['name']} by {', '.join(a['name'] for a in track['artists'])}"}
    elif action == "pause":
        sp.pause_playback(device_id=dev_id); return {"message": "⏸ Paused"}
    elif action == "play":
        sp.start_playback(device_id=dev_id); return {"message": "▶ Resumed"}
    elif action == "next":
        sp.next_track(device_id=dev_id); return {"message": "⏭ Next track"}
    elif action == "prev":
        sp.previous_track(device_id=dev_id); return {"message": "⏮ Previous track"}
    elif action == "volume" and volume is not None:
        sp.volume(int(volume), device_id=dev_id); return {"message": f"🔊 Volume {volume}%"}
    elif action == "status":
        cur = sp.current_user_playing_track()
        if cur and cur["is_playing"]:
            t = cur["item"]; return {"message": f"▶ {t['name']} by {', '.join(a['name'] for a in t['artists'])}"}
        return {"message": "⏸ Nothing playing"}
    return {"error": f"Unknown action: {action}"}

EXECUTORS = {
    "run_command": exec_run_command,
    "open_app":    exec_open_app,
    "type_text":   exec_type_text,
    "screenshot":  exec_screenshot,
    "create_ppt":  exec_create_ppt,
    "create_doc":  exec_create_doc,
    "web_search":  exec_web_search,
    "spotify":     exec_spotify,
    "custom":      exec_custom,
}


# ── Job processing ────────────────────────────────────────────────────────────
def process_job(pair_id: str, secret: str, job: dict):
    job_id = job["job_id"]
    task   = job["task"]
    params = job.get("params", {})

    print(f"\n  ▶ Job {job_id[:8]}… — {task}")
    print(f"    Params: {json.dumps(params)[:120]}")

    executor = EXECUTORS.get(task)
    if not executor:
        agent_post("/agent/jobs/complete", pair_id, secret, {
            "job_id": job_id, "status": "failed",
            "error": f"Unknown task type: {task}",
        })
        print(f"  ❌ Unknown task: {task}")
        return

    # Mark running + update progress
    agent_post("/agent/jobs/complete", pair_id, secret, {
        "job_id": job_id, "status": "running", "progress": 10,
    })

    try:
        result = executor(params)
        success = "error" not in result
        agent_post("/agent/jobs/complete", pair_id, secret, {
            "job_id": job_id,
            "status": "done" if success else "failed",
            "progress": 100,
            "result": result,
            "error": result.get("error") if not success else None,
        })
        if success:
            print(f"  ✅ Done — {result.get('message', 'OK')}")
        else:
            print(f"  ❌ Error — {result.get('error')}")
    except Exception as e:
        agent_post("/agent/jobs/complete", pair_id, secret, {
            "job_id": job_id, "status": "failed", "progress": 0,
            "error": str(e),
        })
        print(f"  ❌ Exception: {e}")


# ── Main loop ─────────────────────────────────────────────────────────────────
def main():
    print("╔══════════════════════════════════════╗")
    print("║       Orin AI PC Agent v1.0          ║")
    print("╚══════════════════════════════════════╝")
    print(f"  Platform: {platform.system()} {platform.machine()}")
    print(f"  Python:   {sys.version.split()[0]}")
    print()

    state = load_state()

    # ── Pair if no existing state ─────────────────────────────────────────
    if not state.get("pair_id") or not state.get("hmac_secret"):
        print("No pairing found. Let's pair with orinai.org.\n")
        print("1. Go to orinai.org → Agent Mode → Desktop Agent tab")
        print("2. Click 'Generate Code' — you'll see a pair_id and a 6-char code")
        print()
        pair_code = input("Enter the 6-char code from orinai.org: ").strip().upper()
        print()
        try:
            pair_data = do_pair(pair_code)
            state.update(pair_data)
            save_state(state)
            print(f"\n  ✅ Paired! pair_id: {state['pair_id'][:16]}…")
        except Exception as e:
            print(f"\n  ❌ Pairing failed: {e}")
            sys.exit(1)
    else:
        print(f"  Using saved pairing: {state['pair_id'][:16]}…")

    pair_id = state["pair_id"]
    secret  = state["hmac_secret"]

    print(f"\n  Polling for jobs every {POLL_INTERVAL_SEC}s. Press Ctrl+C to stop.\n")
    print("  ─────────────────────────────────────")

    last_ping = 0.0

    while True:
        now = time.time()

        # Keepalive ping
        if now - last_ping >= PING_INTERVAL_SEC:
            try:
                agent_post("/agent/ping", pair_id, secret, {})
                last_ping = now
                print(f"  [{datetime.now().strftime('%H:%M:%S')}] 💓 Ping sent — agent online")
            except Exception as e:
                print(f"  [{datetime.now().strftime('%H:%M:%S')}] ⚠️  Ping failed: {e}")

        # Poll for next job
        try:
            data = agent_post("/agent/jobs/next", pair_id, secret, {"wait": LONG_POLL_WAIT})
            job  = data.get("job")
            if job:
                process_job(pair_id, secret, job)
            else:
                # No job — silent wait
                pass
        except KeyboardInterrupt:
            print("\n\n  Goodbye! Agent stopped.")
            break
        except Exception as e:
            print(f"  [{datetime.now().strftime('%H:%M:%S')}] ⚠️  Poll error: {e}")

        try:
            time.sleep(POLL_INTERVAL_SEC)
        except KeyboardInterrupt:
            print("\n\n  Goodbye! Agent stopped.")
            break


if __name__ == "__main__":
    main()
